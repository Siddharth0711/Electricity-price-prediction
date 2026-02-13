"""
Generate PowerPoint Presentation for Mid-Review
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Weather-Driven Electricity Price Forecasting"
    subtitle.text = "Foundation Project 1 - Mid Review\nFebruary 14, 2026"
    
    return slide


def create_content_slide(prs, title_text, content_points):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = title_text
    
    tf = body.text_frame
    for i, point in enumerate(content_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0


def create_presentation():
    """Create complete 10-slide executive mid-review presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. Title
    create_title_slide(prs)
    
    # 2. Business Problem & Success Criteria
    create_content_slide(prs, "Business Problem & Success Criteria", [
        "Business Problem: Volatile price spikes in Indian Day-Ahead Market (DAM).",
        "Business Objective: Forecast volatile electricity prices reliably.",
        "Success Criteria (Business): >10% improvement in trading profit margins.",
        "Success Criteria (Economic): Sharpe Ratio > 1.5; Annualized ROI > 15%.",
        "Feasibility: ML applicability & legal grid code compliance verified."
    ])
    
    # 3. Data Ingestion & Governance
    create_content_slide(prs, "Data Ingestion & Governance", [
        "Data Requirements: High-fidelity 15-min granularity; missingness < 1%.",
        "Data Version Control: Timestamped snapshots for data lineage.",
        "Data Verification: Cross-validated market reports vs actual dispatch logs.",
        "Description: 15-min Market Clearing Prices (MCP) + Multi-Hub weather indices."
    ])
    
    # 4. Feature Engineering: 150+ Features
    create_content_slide(prs, "Feature Engineering: 15-min Granularity", [
        "Time Features: Cyclical encoding (sin/cos) for 15-min blocks.",
        "Weather-Renewable Mapping: Physics-based solar/wind power curves.",
        "Supply Mix Economics: Marginal cost calculation via Merit Order Dispatch.",
        "Lag & Rolling: 24-block (6-hour) window for historical momentum.",
        "Interactions: Demand × Renewable Penetration interactions."
    ])
    
    # 5. Innovation: Merit Order Dispatch
    create_content_slide(prs, "Innovation: Merit Order Dispatch Economics", [
        "Modeling the Grid Supply Curve (Renewables → Nuclear → Coal → Gas).",
        "MCP = Marginal cost of the last dispatched generator.",
        "Innovation: Weather determines the 'Zero-Fuel' supply shift.",
        "Impact: Physics-informed features improve model interpretability."
    ])
    
    # 6. Model Building: Quality & Performance
    create_content_slide(prs, "Model Building: Quality & Performance", [
        "Quality Measures: Performance, Robustness, Scalability, Explainability.",
        "Complexity: Tuned depth (8 levels) to balance bias-variance tradeoff.",
        "Resource Demand: Optimized <50MB RAM footprint for FastAPI serving.",
        "Implementation: Ensemble (XGBoost, LGBM) + Merit-Order physics."
    ])
    
    # 7. Quality Assurance: Reproducibility
    create_content_slide(prs, "Quality Assurance: Reproducibility", [
        "Method Reproducibility: Seeded random states for temporal splits.",
        "Result Reproducibility: Dockerized environment validation.",
        "Experimental Documentation: All hyperparameters logged in YAML configs.",
        "Model Compression: Saved as optimized joblib artifacts for production."
    ])
    
    # 8. Deployment: Production Strategy
    create_content_slide(prs, "Deployment: Production Strategy", [
        "Hardware: Ready for CPU-optimized cloud instances (AWS t3.medium).",
        "Strategy: Blue-Green deployment with Baseline model fallback.",
        "User Acceptance: Simulation on historical hold-out trading sets.",
        "Inference: REST API serving <100ms latency rounds."
    ])
    
    # 9. Monitoring: Maintenance & Drift
    create_content_slide(prs, "Monitoring: Maintenance & Drift", [
        "Non-Stationary Data: Drift monitoring (KS-Test) on distribution shifts.",
        "System Health: Real-time tracking of hardware & infrastructure fatigue.",
        "Maintenance: Automated CI/CD triggers and system update loops.",
        "Quality Logic: Continuous feedback logging & performance triggers."
    ])
    
    # 10. Conclusion & Next Steps
    create_content_slide(prs, "Conclusion & Next Steps", [
        "Current Status: Phase 2 (Real Data Integration) COMPLETE.",
        "Key Deliverables: 1,630+ lines of Python, CRISP-ML(Q) Docs.",
        "Next Steps:",
        "  • Incorporate POSOCO live demand scraping.",
        "  • Refine trading position sizing with transaction costs.",
        "  • Final stakeholders review (Feb 24)."
    ])
    
    return prs

if __name__ == "__main__":
    print("Creating 10-slide executive presentation...")
    prs = create_presentation()
    output_path = "MidReview_ExecutiveSummary.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
