"""
Generate a High-Impact, Professor-Grade Mid-Review Presentation
Focus: Rigorous methodology, crisp architecture, and business-economic impact.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Weather-Driven Electricity Price Forecasting"
    subtitle.text = "Strategic AI for Energy Markets | Foundation Project 1 - Mid Review\nAcademic Auditor: [Professor Name] | Lead: Babita Kiron Vedantam"
    
    # Styling for an 'academic' look (subtle and professional)
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True

def create_content_slide(prs, title_text, items, is_table=False):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(item, tuple):
            p.text = item[0]
            p.level = item[1]
        else:
            p.text = item
            p.level = 0
            
    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    # 1. Title
    create_title_slide(prs)
    
    # 2. Executive Rationale & Business Understanding
    create_content_slide(prs, "1. Executive Rationale: The Business Case", [
        "The Problem: Inefficiency in Indian Day-Ahead Market (DAM) due to high-volatility price spikes.",
        "Primary Objective: Development of a forecast-to-trade pipeline for 15-minute granularity Market Clearing Prices (MCP).",
        "Economic Metrics: Target Sharpe Ratio > 1.5 and an Annualized ROI > 15%.",
        "Technical Threshold: MAPE < 10% on highly volatile peak-load intervals.",
        "Strategic Impact: Reduction of financial risk for portfolio managers via physics-informed forecasting."
    ])
    
    # 3. Methodological Framework: CRISP-ML(Q) Rigor
    create_content_slide(prs, "2. Methodology: Adherence to CRISP-ML(Q)", [
        "Quality-First Lifecycle: Ensuring long-term reliability of mission-critical AI.",
        "Phases Integrated:",
        ("Domain Understanding: Mapping Merit-Order economics to price formation.", 1),
        ("Governed Prep: Data lineage via timestamped versioning.", 1),
        ("Quality Assurance: Reproducibility protocols (Method & Result).", 1),
        ("Production Monitoring: Statistical drift & hardware health metrics.", 1),
        "Compliance: Fully audited against expanded standard criteria."
    ])
    
    # 4. Data Understanding & Governance
    create_content_slide(prs, "3. Data Understanding & Governance", [
        "Granularity: Standardized 15-minute block buckets (96 blocks per Indian trading day).",
        "Multi-Hub Weather Indices: Real-world data from Bhadla (Solar), Muppandal (Wind), and NCR (Demand).",
        "Data Verification: Cross-referencing IEX Market Reports with POSOCO (Grid-India) logs.",
        "Governance Logic: Automatic schema validation & sensor-error imputation (Time-based interpolation).",
        "Lineage: Robust version control for every transformation step."
    ])
    
    # 5. Innovation: Physics-Informed Market Modeling
    create_content_slide(prs, "4. Innovation: Physics-Informed Merit Order", [
        "Market Dynamics Modeling: Constructing the 'Zero-Fuel' Supply Horizon.",
        "Economic Logic: MCP is the marginal cost of the last dispatched generator (Renewables → Nuclear → Coal → Gas).",
        "The Shift: Weather events (solar irradiance/wind speed) shift the supply curve in real-time.",
        "Predictive Edge: Using Merit Order variables instead of pure historical lags significantly reduces bias."
    ])
    
    # 6. Model Architecture: Advanced Ensemble Design
    create_content_slide(prs, "5. Model Architecture: Scalable Ensemble Layer", [
        "Unified Pipeline: Collector → Preprocessor → Feature Builder → Inference Layer.",
        "Algorithmic Stack:",
        ("Baseline: Linear Regression for interpretability benchmarking.", 1),
        ("Advanced: XGBoost, LightGBM, and CatBoost Ensembles.", 1),
        "Feature Density: 146 engineered features including cyclical temporal encoding & interaction terms.",
        "Optimization: Model artifacts compressed for low-latency production serving."
    ])
    
    # 7. Evaluation Strategy: Rigorous Benchmarking
    create_content_slide(prs, "6. Evaluation: Beyond Technical Metrics", [
        "Primary Benchmarks:",
        ("Technical: RMSE 0.047 | MAE 0.033 | MAPE 9%.", 1),
        ("Reproducibility: Seeded validation ensuring consistent results.", 1),
        "Stress Testing: Robustness evaluation during simulated extreme weather heatwaves.",
        "Explainability: SHAP/Feature Importance utilized for 'Physically Interpretable' results.",
        "Bias Recovery: Evaluation of signalWin/Loss ratios in high-volatility spikes."
    ])
    
    # 8. Production Deployment Strategy
    create_content_slide(prs, "7. Deployment: Production-Ready Blueprint", [
        "Architecture: Microservice-ready via FastAPI REST Endpoint.",
        "Performance: Under 100ms inference latency for real-time bidding.",
        "Reliability: Automatic fallback to baseline models during infrastructure anomalies.",
        "Operational Security: Adherence to energy grid data privacy and legal standards.",
        "Hosting: Optimized for CPU-heavy cloud compute instances (AWS/Azure)."
    ])
    
    # 9. Monitoring & Maintenance: Long-Term Reliability
    create_content_slide(prs, "8. Monitoring: Strategic Drift & Health", [
        "Statistical Guardrails: Kolmogorov-Smirnov (KS) tests for feature distribution shifts.",
        "Drift Monitoring: Real-time tracking of non-stationary Indian market demand shifts.",
        "Infrastructure Health: Monitoring hardware fatigue and API response times.",
        "Lifecycle Management: Automated retraining triggers linked to performance decay thresholds."
    ])
    
    # 10. Key Deliverables: Foundation Project Assets
    create_content_slide(prs, "9. Deliverables: Operational Framework", [
        "Production Ecosystem: Collector, Feature Builder, Training, and Server modules.",
        "Documentation: Audited CRISP-ML(Q) Lifecycle Report & Walkthrough.",
        "Market Assets: High-fidelity processed datasets with 96-block daily granularity.",
        "Governance: YAML-based configuration registry for all model hyper-parameters.",
        "API Ready: Fully documented REST end-points for external integration."
    ])
    
    # 11. Timeline & Road Map
    create_content_slide(prs, "10. Status & Future Trajectory", [
        "Current Phase: Mid-Review Audit | Integration Complete.",
        "Milestones Achieved: Real Data Ingestion, 15-min Refactoring, API Setup.",
        "Upcoming Roadmap:",
        ("Phase 4: Real-time scraping from POSOCO/Grid-India logs.", 1),
        ("Phase 5: Refinement of trading position sizing with transaction costs.", 1),
        ("Final Review: February 24, 2026.", 1)
    ])
    
    # 12. Conclusion & Discussion
    create_content_slide(prs, "11. Synthesis & Strategic Outlook", [
        "Conclusion: The project demonstrates a robust, physics-informed AI transition from experimental code to enterprise-grade forecasting.",
        "Key Differentiator: Blending Merit Order economics with high-performance gradient boosting.",
        "Auditor Feedback & Discussion.",
        "Thank You."
    ])

    return prs

if __name__ == "__main__":
    print("Preparing Professor-Grade Academic Presentation...")
    prs = create_presentation()
    output_path = "MidReview_Professor_Audit_Deck.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
